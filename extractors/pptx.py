from pathlib import Path

from pptx import Presentation


def extract_pptx(path: Path, max_chars: int) -> str:
    prs = Presentation(str(path))
    slide_count = len(prs.slides)

    parts = [f"[PPTX] {path.name} - 총 {slide_count}슬라이드"]
    char_count = 0
    truncated = False

    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slide_text = f"\n--- Slide {idx} ---\n" + "\n".join(texts)
        if notes:
            slide_text += f"\n[Notes] {notes}"

        if char_count + len(slide_text) > max_chars:
            truncated = True
            break

        parts.append(slide_text)
        char_count += len(slide_text)

    if truncated:
        parts.append(f"\n[문자 수 제한({max_chars})으로 이후 슬라이드 생략]")

    return "\n".join(parts)
